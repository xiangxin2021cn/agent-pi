# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx>=1.0,<2", "click>=8.3,<9"]
# ///
"""PowerPoint (.pptx) operations tool.

Commands: create, info, extract.

Usage:
    uv run pptx_tool.py COMMAND [OPTIONS]
"""

import json
import re
import sys
from pathlib import Path

import click
from pptx import Presentation


def write_output(text: str, output_path: str | None) -> None:
    """Write text to file or stdout."""
    if output_path:
        Path(output_path).write_text(text, encoding="utf-8")
        click.echo(f"Output written to {output_path}", err=True)
    else:
        click.echo(text)


@click.group()
def cli() -> None:
    """PowerPoint (.pptx) operations tool."""
    pass


@cli.command()
@click.option("--from-file", type=click.Path(exists=True, dir_okay=False), default=None, help="Input markdown or JSON file.")
@click.option("--text", type=str, default=None, help="Inline markdown/text content for slides.")
@click.option("--json-data", type=str, default=None, help="JSON string defining slides structure.")
@click.option("--title", type=str, default=None, help="Presentation title (creates title slide).")
@click.option("--template", type=click.Path(exists=True, dir_okay=False), default=None, help="Template .pptx file to use.")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output .pptx file path.")
def create(from_file: str | None, text: str | None, json_data: str | None, title: str | None, template: str | None, output: str) -> None:
    """Create a PowerPoint presentation from markdown, text, or JSON.

    Markdown format: Use '---' to separate slides. '#' for titles, body text below.

    JSON format:
    [
      {"title": "Slide Title", "body": "Bullet 1\\nBullet 2", "notes": "Speaker notes"},
      {"title": "Another Slide", "body": "Content here"}
    ]
    """
    if from_file is None and text is None and json_data is None and title is None:
        click.echo("Error: provide --from-file, --text, --json-data, or --title.", err=True)
        sys.exit(1)

    try:
        if template:
            prs = Presentation(template)
        else:
            prs = Presentation()

        slides_data: list[dict[str, str]] = []

        if json_data:
            # Parse JSON
            data_path = Path(json_data)
            if data_path.exists() and data_path.is_file():
                parsed = json.loads(data_path.read_text(encoding="utf-8"))
            else:
                parsed = json.loads(json_data)

            if isinstance(parsed, list):
                slides_data = parsed
            else:
                click.echo("Error: JSON must be an array of slide objects.", err=True)
                sys.exit(1)

        elif from_file:
            content = Path(from_file).read_text(encoding="utf-8")
            if from_file.endswith(".json"):
                parsed = json.loads(content)
                if isinstance(parsed, list):
                    slides_data = parsed
                else:
                    click.echo("Error: JSON must be an array of slide objects.", err=True)
                    sys.exit(1)
            else:
                slides_data = _parse_markdown_slides(content)

        elif text:
            slides_data = _parse_markdown_slides(text)

        # Add title slide if requested
        if title:
            _add_title_slide(prs, title)

        # Add content slides
        for slide_info in slides_data:
            _add_content_slide(prs, slide_info)

        prs.save(output)
        click.echo(f"Presentation saved to {output} ({len(slides_data) + (1 if title else 0)} slides)", err=True)
    except json.JSONDecodeError as e:
        click.echo(f"Error parsing JSON: {e}", err=True)
        sys.exit(1)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


def _parse_markdown_slides(md_text: str) -> list[dict[str, str]]:
    """Parse markdown into slide data. Slides separated by '---'."""
    raw_slides = re.split(r"\n---\n|\n---$|^---\n", md_text)
    slides: list[dict[str, str]] = []

    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue

        lines = raw.split("\n")
        slide_title = ""
        body_lines: list[str] = []

        for line in lines:
            heading_match = re.match(r"^#{1,3}\s+(.+)$", line.strip())
            if heading_match and not slide_title:
                slide_title = heading_match.group(1)
            else:
                body_lines.append(line)

        body = "\n".join(body_lines).strip()
        slide: dict[str, str] = {}
        if slide_title:
            slide["title"] = slide_title
        if body:
            slide["body"] = body
        if slide:
            slides.append(slide)

    return slides


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    """Add a title slide."""
    try:
        layout = prs.slide_layouts[0]  # Title Slide layout
    except IndexError:
        click.echo("Error: template has no slide layouts.", err=True)
        sys.exit(1)
    slide = prs.slides.add_slide(layout)

    if slide.placeholders:
        # Title placeholder
        if 0 in slide.placeholders:
            slide.placeholders[0].text = title
        # Subtitle placeholder
        if 1 in slide.placeholders and subtitle:
            slide.placeholders[1].text = subtitle


def _add_content_slide(prs: Presentation, slide_info: dict[str, str]) -> None:
    """Add a content slide from slide info dict."""
    title = slide_info.get("title", "")
    body = slide_info.get("body", "")
    notes = slide_info.get("notes", "")

    # Use Title and Content layout (index 1) or Blank (index 6)
    try:
        if title and body:
            layout = prs.slide_layouts[1]  # Title and Content
        elif title:
            layout = prs.slide_layouts[5]  # Title Only
        else:
            layout = prs.slide_layouts[6]  # Blank
    except IndexError:
        if not prs.slide_layouts:
            click.echo("Error: template has no slide layouts.", err=True)
            sys.exit(1)
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)

    # Set title
    if title and slide.placeholders and 0 in slide.placeholders:
        slide.placeholders[0].text = title

    # Set body content
    if body:
        if slide.placeholders and 1 in slide.placeholders:
            tf = slide.placeholders[1].text_frame
        else:
            # No content placeholder (e.g. Blank layout) — add a textbox
            from pptx.util import Inches
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            tf = txBox.text_frame
            tf.word_wrap = True
        tf.clear()

        body_lines = body.split("\n")
        first_line = True
        for line in body_lines:
            stripped = line.strip()
            if not stripped:
                continue

            if first_line:
                p = tf.paragraphs[0]
                first_line = False
            else:
                p = tf.add_paragraph()

            # Handle bullet points
            bullet_match = re.match(r"^[-*+]\s+(.+)$", stripped)
            num_match = re.match(r"^\d+[.)]\s+(.+)$", stripped)

            if bullet_match:
                p.text = bullet_match.group(1)
                p.level = 0
            elif num_match:
                p.text = num_match.group(1)
                p.level = 0
            else:
                # Check for indented bullets or numbered sub-items
                indent_match = re.match(r"^(\s+)[-*+]\s+(.+)$", line)
                indent_num_match = re.match(r"^(\s+)\d+[.)]\s+(.+)$", line)
                if indent_match:
                    indent_level = len(indent_match.group(1)) // 2
                    p.text = indent_match.group(2)
                    p.level = min(indent_level, 4)
                elif indent_num_match:
                    indent_level = len(indent_num_match.group(1)) // 2
                    p.text = indent_num_match.group(2)
                    p.level = min(indent_level, 4)
                else:
                    p.text = stripped

    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("-o", "--output", type=click.Path(), default=None, help="Write output to file.")
def info(file: str, output: str | None) -> None:
    """Show presentation metadata and slide information."""
    try:
        prs = Presentation(file)

        slides_info = []
        for i, slide in enumerate(prs.slides):
            slide_data: dict[str, object] = {
                "number": i + 1,
                "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
            }

            # Get title
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text

            # Count shapes by type
            slide_data["shape_count"] = len(slide.shapes)

            # Notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_data["has_notes"] = True

            slides_info.append(slide_data)

        # Slide dimensions
        width_in = prs.slide_width / 914400  # EMU to inches
        height_in = prs.slide_height / 914400

        info_dict: dict[str, object] = {
            "file": str(Path(file).resolve()),
            "slide_count": len(prs.slides),
            "slide_width_inches": round(width_in, 2),
            "slide_height_inches": round(height_in, 2),
            "slides": slides_info,
        }

        # Available layouts
        layouts = [layout.name for layout in prs.slide_layouts]
        info_dict["available_layouts"] = layouts

        write_output(json.dumps(info_dict, indent=2, default=str), output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--slide", type=int, default=None, help="Extract specific slide number (1-based).")
@click.option("--include-notes/--no-notes", default=True, help="Include speaker notes (default: yes).")
@click.option("-o", "--output", type=click.Path(), default=None, help="Write output to file.")
def extract(file: str, slide: int | None, include_notes: bool, output: str | None) -> None:
    """Extract text content from a PowerPoint presentation."""
    try:
        prs = Presentation(file)
        parts: list[str] = []

        slides_to_process = []
        if slide is not None:
            if 1 <= slide <= len(prs.slides):
                slides_to_process = [(slide - 1, prs.slides[slide - 1])]
            else:
                click.echo(f"Error: slide {slide} out of range (1-{len(prs.slides)}).", err=True)
                sys.exit(1)
        else:
            slides_to_process = list(enumerate(prs.slides))

        for i, s in slides_to_process:
            parts.append(f"--- Slide {i + 1} ---")

            # Title
            if s.shapes.title:
                parts.append(f"# {s.shapes.title.text}")

            # All text shapes
            for shape in s.shapes:
                if shape.has_text_frame:
                    # Skip title shape (already handled)
                    if shape == s.shapes.title:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            # Add indentation for bullet levels
                            indent = "  " * paragraph.level if paragraph.level else ""
                            parts.append(f"{indent}- {text}" if paragraph.level > 0 else text)

                # Table content
                if shape.has_table:
                    table = shape.table
                    parts.append("")
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        parts.append(" | ".join(cells))

            # Speaker notes
            if include_notes and s.has_notes_slide:
                notes_text = s.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"\n[Notes: {notes_text}]")

            parts.append("")

        write_output("\n".join(parts), output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


if __name__ == "__main__":
    cli()
